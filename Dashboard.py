# Import necessary libraries
import streamlit as st
import pandas as pd
import os
import plotly.express as px
from geopy.geocoders import Nominatim
import plotly.graph_objects as go
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import plotly.io as pio
import datetime

# Set page configuration
st.set_page_config(layout="wide")
TIME_STAMP = datetime.datetime.now()

# Define CSS styles
metric_style = """
<style>
.metric-container {
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    height: 150px;
    width: 180px;
    border: 1px solid #ccc;
    border-radius: 5px;
    padding: 10px;
    box-sizing: border-box;
    margin-top:60px;
    margin-bottom:100px;
    background-color: RGB(245,245,245)

}
.metric-label {
    font-size: 16px;
    padding-top:10px
}
.metric-value {
    font-size: 70px;
    color: #008000;
}
</style>
"""

st.markdown(metric_style, unsafe_allow_html=True)

# Define constants
EXCEL_FILE_WITH_COORDINATES = "./data/output_excel.xlsx"
REPORT_URL = 'https://elekta.my.salesforce.com/00O6g000006RqPX'
MAP_PATH = 'graphs/map.png'

# Initialize dataframe
df = pd.DataFrame()

# Set page title
st.title('Region Overview')

# Sidebar for data loading
st.sidebar.title('Load Data')
uploaded_file = st.sidebar.file_uploader("Choose CSV File", type=['csv'])

# Check if a file is uploaded
if uploaded_file:
    # Load data from existing excel file or uploaded csv file
    # if os.path.exists(EXCEL_FILE_WITH_COORDINATES):
    #     df = pd.read_excel(EXCEL_FILE_WITH_COORDINATES)
    # else:
    df = pd.read_csv(uploaded_file)

    # Rename columns for consistency
    df.rename(columns={
        'Account': 'account',
        'Location': 'location',
        'IP Street': 'street',
        'IP City': 'city',
        'IP State': 'state',
        'IP Zip/Postal Code': 'zipcode',
        'Primary Technician: Member Name': 'primary_fse',
        'Secondary Technician Name': 'secondary_fse',
        'EoL Date IP': 'eol',
        'Device Age': 'device_age',
        'Customer/Device Acceptance Date': 'cat',
        'EoGS Date IP': 'eogs',
        'Installed Product: Installed Product': 'ip',
        'Primary Technician: City': 'technician_city',
        'Primary Technician: Zip': 'technician_zip',
        'Primary Technician: State': 'technician_state',
        'Primary Technician: Street': 'technician_street',
        'Primary Technician: Service Manager': 'manager'

    }, inplace=True)

  
    # Clean data and add address columns for IPs and technicians
    df_clean = df.dropna(subset=['location'])
    df_clean['address'] = df_clean.apply(lambda row: f"{row['street']}, {row['city']}, {row['state']}, {str(row['zipcode'])[:5]}", axis=1)

    # Initialize geocoder
    geolocator = Nominatim(user_agent="LocatingMachines")
    def geocode_address(address):
        print(f"Geocoding: {address}")
        try:
            location = geolocator.geocode(address, timeout=10)

            if location:
                print(f"Found: {location.latitude}, {location.longitude}")  # Print the found coordinates
                return pd.Series([location.latitude, location.longitude])
            else:
                print("Location not found!")
                return pd.Series([None, None])

        except Exception as e:
            print(f"Error: {e}")
            return pd.Series([None, None])
        


    # Function to get latitude and longitude from an address
    manager = df_clean.loc[2, 'manager']
    df_clean[['latitude', 'longitude']] = df_clean['address'].apply(geocode_address)
    directory = f'data/{manager}'
    if not os.path(directory):
        os.makedirs(directory)
   
    df_clean.to_excel(f'./data/{manager}/report_{manager}_{TIME_STAMP}.xlsx', sheet_name='data')
    df_clean_no_coordinates = df_clean.dropna(subset=['latitude', 'longitude'])
    df_clean_no_coordinates.to_excel(f'./data/{manager}report_{manager}_{TIME_STAMP}.xlsx', sheet_name='data')


    # Display region metrics
    container = st.container(border=True)
    head_count = len(df_clean['primary_fse'].unique())
    average_equipment_age = df_clean['device_age'].mean()
    device_type_frequency = df_clean['ip'].str.split('/').str[0].value_counts()
    with container:
        col1, col2 = st.columns(2)
        with col1:
            cola, colb,colc = st.columns(3)
            with cola:
                st.markdown(f"<div class='metric-container'><div class='metric-label'>Head Count</div><div class='metric-value'>{head_count}</div></div>", unsafe_allow_html=True)
            with colb:
                st.markdown(f"<div class='metric-container'><div class='metric-label'>Average IP Age</div><div class='metric-value'>{round(average_equipment_age, 2)}</div></div>", unsafe_allow_html=True)
            with colc:
                st.markdown(f"<div class='metric-container'><div class='metric-label'>Total IPs</div><div class='metric-value'> {len(df_clean['device_age'])}</div></div>", unsafe_allow_html=True)

            st.divider()
            for i in range(0, len(device_type_frequency), 3):
                cols = st.columns(3)
                for j in range(3):
                    if i + j < len(device_type_frequency):
                        cols[j].metric(label=device_type_frequency.index[i + j], value=device_type_frequency.values[i + j])
        with col2:
            # Plot device locations on a map
            # Step 1: Find unique FSEs and corresponding address
            df_clean['technician_address'] = df_clean.apply(lambda row: f"{row['technician_street']},{row['technician_city']}, {row['technician_zip']}, {str(row['technician_state'])[:5]}", axis=1)
            
            unique_fse_df = df_clean.drop_duplicates(subset='primary_fse')
            unique_fse_df[['latitude', 'longitude']] = unique_fse_df['technician_address'].apply(geocode_address)

            # Step 2: Combine technician_city, technician_zip, and technician_state into one 'technician_address' column
            # Step 1: Add a 'dataset' column to differentiate the two sets of data
            df_clean['dataset'] = 'IPs'
            unique_fse_df['dataset'] = 'FSEs'
            # Step 2: Combine both datasets
            combined_df = pd.concat([df_clean, unique_fse_df], ignore_index=True)

            # Make the size column for larger dots
            combined_df['size'] = [15 if dataset == 'unique_fse' else 10 for dataset in combined_df['dataset']]


            map_fig = px.scatter_mapbox(combined_df, 
                                    lat="latitude", 
                                    lon="longitude", 
                                    zoom=4, 
                                    height=750, 
                                    width=750, 
                                    color='dataset', 
                                    hover_name='primary_fse', 
                                    hover_data=['account'], 
                                    size='size',
                                    size_max=10,  # Maximum size of dots
                                    color_discrete_map={
                                    'IPs': 'green',  # Color for original data points
                                    'FSEs': 'red'  # Color for unique FSE data points
                                    },)
            
            map_fig.update_layout(mapbox_style="carto-positron")
            map_fig.update_layout(margin={"r":0,"t":0,"l":0,"b":0})
            
            st.plotly_chart(map_fig)
            pio.write_image(map_fig, MAP_PATH)



    # Plot device age
    st.title('IP Age')
    sorted_data = df_clean.sort_values(by='device_age')
    colors = ['red' if age > 12 else 'blue' for age in sorted_data['device_age']]
    device_age_fig = go.Figure(go.Bar(y=sorted_data['device_age'], x=sorted_data['ip'],  marker=dict(color=colors)))
    device_age_fig.add_hline(y=12, line_width=3, line_dash="dash", line_color="gray", annotation_text="Due for replacement age > 12", annotation_position="top left")
    device_age_fig.update_layout(height=700)
    st.plotly_chart(device_age_fig)
    device_age_fig.write_image('graphs/device_age_fig.png')

    # Plot EOL timelines
    st.title('IP EOL Timeline')
    df_clean['eol'] = pd.to_datetime(df_clean['eol'], errors='coerce')
    df_clean['eogs'] = pd.to_datetime(df_clean['eogs'], errors='coerce')
    today = datetime.today()

    df_filtered = df_clean.dropna(subset=['eol', 'eogs'])
    df_filtered['eogs'] = pd.to_datetime(df_filtered['eogs'], errors='coerce')
    df_timeline = pd.DataFrame({'Account': df_filtered['account'], 'Product': df_filtered['ip'], 'Start': [today] * len(df_filtered), 'End': df_filtered['eol']})
    df_timeline['Status'] = df_timeline['End'].apply(lambda x: 'Past' if x < today else 'Future')
    fig = px.timeline(df_timeline, x_start='Start', x_end='End', y='Product', height=1000, color='Status', color_discrete_map={'Past': 'red', 'Future': 'green'}, title="EOL Product Timeline")
    fig.update_layout(xaxis_title="Date", yaxis_title="Products", xaxis=dict(tickformat="%m-%d-%Y"), showlegend=False)

    st.plotly_chart(fig)

    # Plot EOGS timeline
    st.title('IP EOGS Timeline')
    df_timeline = pd.DataFrame({'Account': df_filtered['account'], 'Product': df_filtered['ip'], 'Start': [today] * len(df_filtered), 'End': df_filtered['eogs']})
    df_timeline['Status'] = df_timeline['End'].apply(lambda x: 'Past' if x < today else 'Future')
    fig = px.timeline(df_timeline, x_start='Start', x_end='End', y='Product', height=1000, color='Status', color_discrete_map={'Past': 'red', 'Future': 'green'}, title="EOGS Product Timeline")
    fig.update_layout(xaxis_title="Date", yaxis_title="Products", xaxis=dict(tickformat="%m-%d-%Y"), showlegend=False)
    st.plotly_chart(fig)

        # Create a new presentation
    prs = Presentation()

    # Ensure the 'graphs' directory exists
    if not os.path.exists('graphs'):
        os.makedirs('graphs')

    # Slide 1: Main Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title = slide.shapes.title
    title.text = "Region Overview"

    # Add head_count, average_equipment_age, and device_type_frequency as text
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"Head Count: {head_count}\nAverage IP Age: {round(average_equipment_age, 2)}\nDevice Type Frequency: {device_type_frequency}"

    # Save map as image and add to slide
    map_path = 'graphs/map.png'
    fig.write_image(map_path)

    left = Inches(5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(5)

    pic = slide.shapes.add_picture(map_path, left, top, width, height)

    # Slide 2: Device Age
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title = slide.shapes.title
    title.text = "Device Age"

    # Save device_age_fig as image and add to slide
    device_age_fig_path = 'graphs/device_age_fig.png'
    device_age_fig.write_image(device_age_fig_path)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    pic = slide.shapes.add_picture(device_age_fig_path, left, top, width, height)

    # Save the presentation
    prs.save('region_overview.pptx')


# Display instructions if no file is uploaded
else:
    st.header('Upload Report')
    st.write('1. CLM Primary Sites Filter by Manger')
    st.markdown(f"Report on CLM({REPORT_URL}): ")
    st.write('2. Customized and enter RSM name.')
    st.write('3. Run report')
    st.write('4. Export to csv file.')


